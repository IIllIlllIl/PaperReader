"""
Tests for PPT Generator
"""

from base64 import b64decode

import pytest
from pptx import Presentation as PPTXPresentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from src.generation.ppt_generator import PPTGenerator
from src.generation.pptx_exporter import markdown_to_pptx
from src.analysis.content_extractor import ContentExtractor, OrganizedPresentation, SlideContent


class TestPPTGenerator:
    """Test PPT generator"""

    def test_generator_init(self):
        """Test generator initialization"""
        generator = PPTGenerator()
        assert generator is not None

    def test_generate_markdown(self):
        """Test Markdown generation"""
        generator = PPTGenerator()

        # Create sample presentation
        slides = [
            SlideContent(title="Test Slide", bullet_points=["Point 1", "Point 2"], notes="Test notes")
        ]
        presentation = OrganizedPresentation(slides=slides, total_slides=1)

        # Generate Markdown
        markdown = generator.generate_markdown(presentation)

        assert markdown is not None
        assert len(markdown) > 0
        assert "marp: true" in markdown
        assert "Test Slide" in markdown

    def test_generate_markdown_preserves_figure_markdown(self):
        """Test figure slides keep image markdown and captions"""
        generator = PPTGenerator()

        slides = [
            SlideContent(
                title="Figure Slide",
                bullet_points=[
                    "![Figure Slide](outputs/images/example.png)",
                    "*Example caption*"
                ],
                slide_type="figure"
            )
        ]
        presentation = OrganizedPresentation(slides=slides, total_slides=1)

        markdown = generator.generate_markdown(presentation)

        assert "![Figure Slide](outputs/images/example.png)" in markdown
        assert "*Example caption*" in markdown
        assert "- ![Figure Slide](outputs/images/example.png)" not in markdown

    def test_generate_markdown_adds_figure_fields(self):
        """Test PhD-style figure metadata is bridged into Markdown"""
        generator = PPTGenerator()

        class FigureSlide:
            title = "Method Overview"
            bullet_points = ["Key idea", "Component A"]
            notes = "Overview"
            slide_type = "content"
            figure_path = "outputs/images/overview.png"
            figure_caption = "System architecture"
            has_figure = True

        presentation = OrganizedPresentation(slides=[FigureSlide()], total_slides=1)

        markdown = generator.generate_markdown(presentation)

        assert "- Key idea" in markdown
        assert "- Component A" in markdown
        assert "![Method Overview](outputs/images/overview.png)" in markdown
        assert "*System architecture*" in markdown

    def test_markdown_to_pptx_renders_image(self, tmp_path):
        """Test PPTX exporter renders Markdown image syntax"""
        image_path = tmp_path / "figure.png"
        image_path.write_bytes(b64decode(
            "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAusB9Wn9lfsAAAAASUVORK5CYII="
        ))

        markdown_path = tmp_path / "figure.md"
        markdown_path.write_text(
            "---\n"
            "marp: true\n"
            "---\n\n"
            "## Figure Slide\n\n"
            f"![Figure Slide]({image_path})\n"
            "*Example caption*\n",
            encoding="utf-8"
        )

        output_path = tmp_path / "figure.pptx"
        markdown_to_pptx(str(markdown_path), str(output_path), title="Figure Deck")

        assert output_path.exists()

        presentation = PPTXPresentation(str(output_path))
        assert len(presentation.slides) == 1
        assert any(
            shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            for shape in presentation.slides[0].shapes
        )

    def test_save_presentation(self, tmp_path):
        """Test saving presentation to file"""
        generator = PPTGenerator()

        markdown = "---\nmarp: true\n---\n\n## Test\n\n- Point 1"
        output_path = tmp_path / "test.md"

        generator.save_presentation(markdown, str(output_path))

        assert output_path.exists()
        assert output_path.read_text() == markdown


class TestContentExtractor:
    """Test content extractor"""

    def test_extractor_init(self):
        """Test extractor initialization"""
        extractor = ContentExtractor()
        assert extractor is not None

    @pytest.mark.skip(reason="V3 API does not have suggest_visualizations method")
    def test_suggest_visualizations(self):
        """Test visualization suggestions (deprecated in V3)"""
        extractor = ContentExtractor()

        results = [
            "Accuracy improved to 95%",
            "Faster than baseline"
        ]

        suggestions = extractor.suggest_visualizations(results)

        assert len(suggestions) <= 5
        assert all('type' in s for s in suggestions)
        assert all('suggestion' in s for s in suggestions)
