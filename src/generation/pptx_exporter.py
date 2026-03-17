#!/usr/bin/env python3
"""
PPTX Exporter

Converts Marp Markdown files to PowerPoint PPTX format
"""

import logging
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re

logger = logging.getLogger(__name__)

IMAGE_MARKDOWN_RE = re.compile(r'^!\[(?P<alt>.*?)\]\((?P<path>.+?)\)$')
ITALIC_CAPTION_RE = re.compile(r'^\*(?P<caption>[^\s].*)\*$')


def parse_marpedown(md_file: str):
    """Parse Marp Markdown file and extract slides"""

    markdown_path = Path(md_file)

    with open(markdown_path, 'r', encoding='utf-8') as f:
        content = f.read()

    # Remove YAML front matter
    if content.startswith('---'):
        parts = content.split('---', 2)
        if len(parts) >= 3:
            content = parts[2]

    # Split by slide separators
    slides_raw = content.split('\n---\n')

    slides = []
    for slide_raw in slides_raw:
        if not slide_raw.strip():
            continue

        slide_data = {
            'title': '',
            'bullets': [],
            'image_path': '',
            'image_alt': '',
            'caption': ''
        }

        for raw_line in slide_raw.strip().split('\n'):
            line = raw_line.strip()
            if not line:
                continue

            if line.startswith('<!--'):
                continue

            image_match = IMAGE_MARKDOWN_RE.match(line)
            caption_match = ITALIC_CAPTION_RE.match(line)

            if line.startswith('## '):
                slide_data['title'] = line[3:].strip()
            elif image_match:
                image_path = Path(image_match.group('path')).expanduser()
                if not image_path.is_absolute():
                    image_path = (markdown_path.parent / image_path).resolve()
                slide_data['image_path'] = str(image_path)
                slide_data['image_alt'] = image_match.group('alt').strip()
            elif caption_match and slide_data['image_path'] and not slide_data['caption']:
                slide_data['caption'] = caption_match.group('caption').strip()
            elif line.startswith('- ') or line.startswith('* '):
                slide_data['bullets'].append(line[2:].strip())

        if any([slide_data['title'], slide_data['bullets'], slide_data['image_path'], slide_data['caption']]):
            slides.append(slide_data)

    return slides


def _add_title(slide, title: str) -> None:
    """Add a styled title to a slide."""
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(12.333)
    height = Inches(1)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title

    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    title_para.alignment = PP_ALIGN.LEFT


def _add_bullets(slide, bullets, left, top, width, height) -> None:
    """Add bullet text to a slide."""
    if not bullets:
        return

    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    for i, bullet in enumerate(bullets[:6]):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = bullet
        p.font.size = Pt(24)
        p.level = 0
        p.space_after = Pt(12)

        if bullet.startswith('✓'):
            p.font.color.rgb = RGBColor(0x27, 0xAE, 0x60)
        elif bullet.startswith('✗'):
            p.font.color.rgb = RGBColor(0xE7, 0x4C, 0x3C)
        else:
            p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


def _add_image(slide, image_path: str, left, top, width, height):
    """Add an image fitted inside a target box."""
    image_file = Path(image_path)
    if not image_file.exists():
        logger.warning(f"Image not found, skipping: {image_path}")
        return None

    picture = slide.shapes.add_picture(str(image_file), left, top, width=width)

    if picture.height > height:
        scale = height / picture.height
        picture.width = int(picture.width * scale)
        picture.height = int(picture.height * scale)

    picture.left = int(left + (width - picture.width) / 2)
    picture.top = int(top + (height - picture.height) / 2)
    return picture


def _add_caption(slide, caption: str, left, top, width, height) -> None:
    """Add caption text below an image."""
    if not caption:
        return

    caption_box = slide.shapes.add_textbox(left, top, width, height)
    caption_frame = caption_box.text_frame
    caption_frame.word_wrap = True
    caption_frame.text = caption

    caption_para = caption_frame.paragraphs[0]
    caption_para.font.size = Pt(14)
    caption_para.font.italic = True
    caption_para.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    caption_para.alignment = PP_ALIGN.CENTER


def _render_content_slide(slide, slide_data) -> None:
    """Render a non-title slide with optional bullets and image."""
    bullets = slide_data['bullets']
    image_path = slide_data['image_path']
    caption = slide_data['caption']

    if image_path and bullets:
        _add_bullets(
            slide,
            bullets,
            Inches(0.75),
            Inches(1.8),
            Inches(5.2),
            Inches(4.9)
        )
        picture = _add_image(
            slide,
            image_path,
            Inches(6.35),
            Inches(1.8),
            Inches(6.2),
            Inches(3.9)
        )
        if picture is not None:
            _add_caption(
                slide,
                caption,
                Inches(6.2),
                Inches(5.95),
                Inches(6.45),
                Inches(0.45)
            )
    elif image_path:
        picture = _add_image(
            slide,
            image_path,
            Inches(0.9),
            Inches(1.7),
            Inches(11.5),
            Inches(4.8)
        )
        if picture is not None:
            _add_caption(
                slide,
                caption,
                Inches(1.0),
                Inches(6.1),
                Inches(11.3),
                Inches(0.45)
            )
    else:
        _add_bullets(
            slide,
            bullets,
            Inches(0.75),
            Inches(2),
            Inches(11.833),
            Inches(5)
        )


def create_pptx(slides, output_file: str, title: str = "Presentation"):
    """Create PPTX from slides"""

    prs = Presentation()

    # Set slide dimensions to 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    start_index = 0

    # Create title slide when the first slide is text-only
    if slides and slides[0]['title'] and not slides[0]['image_path']:
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = slides[0]['title']

        subtitle_lines = slides[0]['bullets'][:2]
        if subtitle_lines:
            subtitle = slide.placeholders[1]
            subtitle.text = '\n'.join(subtitle_lines)

        start_index = 1

    for slide_data in slides[start_index:]:
        if not any([slide_data['title'], slide_data['bullets'], slide_data['image_path'], slide_data['caption']]):
            continue

        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        if slide_data['title']:
            _add_title(slide, slide_data['title'])

        _render_content_slide(slide, slide_data)

    prs.save(output_file)
    logger.info(f"✅ Created PPTX: {output_file}")
    logger.info(f"   Total slides: {len(prs.slides)}")


def markdown_to_pptx(markdown_path: str, output_path: str, title: str = None) -> None:
    """
    Convert Markdown file to PPTX

    Args:
        markdown_path: Path to Markdown file
        output_path: Path to output PPTX file
        title: Presentation title (default: filename)
    """
    if title is None:
        title = Path(markdown_path).stem

    slides = parse_marpedown(markdown_path)
    create_pptx(slides, output_path, title=title)
