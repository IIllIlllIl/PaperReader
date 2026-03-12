#!/usr/bin/env python3
"""
Enhanced Markdown to PPTX Converter - WITH IMAGE SUPPORT

This is a prototype showing the key changes needed to support images.
Compare with the original md_to_pptx.py to see the differences.
"""

import sys
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def parse_marpedown_with_images(md_file: str):
    """
    Parse Marp Markdown file and extract slides WITH IMAGE SUPPORT

    KEY CHANGE: Detect and parse image markdown syntax ![alt](path)
    """
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()

    # Remove YAML front matter
    if content.startswith('---'):
        parts = content.split('---', 2)
        if len(parts) >= 3:
            content = parts[2]

    # Split by slide separators
    slides_raw = content.split('\n---\n')

    slides = []
    for slide_raw in slides_raw:
        if not slide_raw.strip():
            continue

        lines = slide_raw.strip().split('\n')
        title = ""
        bullets = []

        for line in lines:
            line = line.strip()
            if not line:
                continue

            # Skip comments
            if line.startswith('<!--'):
                continue

            # Extract title (## Header)
            if line.startswith('## '):
                title = line[3:].strip()

            # Extract bullet points WITH IMAGE DETECTION ⭐
            elif line.startswith('- ') or line.startswith('* '):
                content_text = line[2:].strip()

                # ⭐ KEY CHANGE: Check if content is an image markdown
                img_match = re.match(r'!\[([^\]]*)\]\(([^\)]+)\)', content_text)

                if img_match:
                    # This is an image reference!
                    bullets.append({
                        'type': 'image',
                        'alt': img_match.group(1),
                        'path': img_match.group(2)
                    })
                    print(f"   🖼️  Found image: {img_match.group(1)} -> {img_match.group(2)}")
                else:
                    # Regular text
                    bullets.append({
                        'type': 'text',
                        'content': content_text
                    })

        if title or bullets:
            slides.append({
                'title': title,
                'bullets': bullets
            })

    return slides


def add_image_to_slide(slide, image_path: str, left: float, top: float,
                      max_width: float, max_height: float):
    """
    Add image to slide with proper sizing

    ⭐ NEW FUNCTION: Handles image insertion logic

    Args:
        slide: python-pptx slide object
        image_path: Path to image file
        left: Left position in inches
        top: Top position in inches
        max_width: Maximum width in inches
        max_height: Maximum height in inches

    Returns:
        True if successful, False otherwise
    """
    # Resolve image path
    img_path = Path(image_path)
    if not img_path.exists():
        print(f"   ⚠️  Image not found: {image_path}")
        return False

    try:
        from PIL import Image

        # Get image dimensions to maintain aspect ratio
        with Image.open(img_path) as img:
            img_width, img_height = img.size

        # Calculate aspect ratio
        aspect = img_width / img_height if img_height > 0 else 1

        # Calculate display size (fit within max bounds)
        if max_width / aspect <= max_height:
            # Width-constrained
            display_width = max_width
            display_height = max_width / aspect
        else:
            # Height-constrained
            display_height = max_height
            display_width = max_height * aspect

        # Add picture to slide
        slide.shapes.add_picture(
            str(img_path),
            Inches(left),
            Inches(top),
            width=Inches(display_width),
            height=Inches(display_height)
        )

        print(f"   ✅ Added image: {img_path.name} ({display_width:.1f}\" x {display_height:.1f}\")")
        return True

    except Exception as e:
        print(f"   ❌ Failed to add image {image_path}: {e}")
        return False


def create_pptx_with_images(slides, output_file: str, title: str = "Presentation"):
    """
    Create PPTX from slides WITH IMAGE SUPPORT

    ⭐ KEY CHANGES:
    1. Detect if slide has images
    2. Create different layouts for text-only vs image slides
    3. Call add_image_to_slide() for images
    """

    prs = Presentation()

    # Set slide dimensions to 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Create title slide
    if slides and slides[0]['title']:
        slide_layout = prs.slide_layouts[0]  # Title slide
        slide = prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = slides[0]['title']

        # Handle subtitle with mixed content
        if slides[0]['bullets']:
            text_items = [b for b in slides[0]['bullets'] if b.get('type') == 'text']
            if text_items:
                subtitle = slide.placeholders[1]
                subtitle.text = '\n'.join([b['content'] for b in text_items[:2]])

    # Create content slides
    for slide_data in slides[1:]:
        if not slide_data['title'] and not slide_data['bullets']:
            continue

        # Use blank layout
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Add title
        if slide_data['title']:
            left = Inches(0.5)
            top = Inches(0.5)
            width = Inches(12.333)
            height = Inches(1)

            title_box = slide.shapes.add_textbox(left, top, width, height)
            title_frame = title_box.text_frame
            title_frame.text = slide_data['title']

            # Style the title
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(36)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
            title_para.alignment = PP_ALIGN.LEFT

        # ⭐ KEY CHANGE: Detect if slide has images
        has_images = any(b.get('type') == 'image' for b in slide_data['bullets'])

        if has_images:
            # Mixed content slide (text + images)
            text_items = [b for b in slide_data['bullets'] if b.get('type') == 'text']
            image_items = [b for b in slide_data['bullets'] if b.get('type') == 'image']

            # Add text on left side (40% width)
            if text_items:
                text_left = Inches(0.75)
                text_top = Inches(2)
                text_width = Inches(4.5)
                text_height = Inches(5)

                text_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
                text_frame = text_box.text_frame
                text_frame.word_wrap = True

                for i, item in enumerate(text_items[:4]):  # Max 4 text items
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()

                    p.text = item['content']
                    p.font.size = Pt(20)
                    p.space_after = Pt(10)
                    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

            # Add image on right side (60% width)
            for i, img_item in enumerate(image_items[:2]):  # Max 2 images
                img_left = 5.5
                img_top = 2.0
                img_max_width = 7.0
                img_max_height = 5.0

                add_image_to_slide(
                    slide,
                    img_item['path'],
                    img_left,
                    img_top,
                    img_max_width,
                    img_max_height
                )

        else:
            # Original text-only logic (unchanged)
            if slide_data['bullets']:
                left = Inches(0.75)
                top = Inches(2)
                width = Inches(11.833)
                height = Inches(5)

                text_box = slide.shapes.add_textbox(left, top, width, height)
                text_frame = text_box.text_frame
                text_frame.word_wrap = True

                for i, bullet in enumerate(slide_data['bullets'][:6]):
                    if bullet.get('type') != 'text':
                        continue

                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()

                    p.text = bullet['content']
                    p.font.size = Pt(24)
                    p.level = 0
                    p.space_after = Pt(12)
                    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # Save presentation
    prs.save(output_file)
    print(f"\n✅ Created PPTX: {output_file}")
    print(f"   Total slides: {len(prs.slides)}")


def main():
    """Main entry point"""
    if len(sys.argv) < 2:
        print("Usage: python md_to_pptx_prototype.py <input.md> [output.pptx]")
        sys.exit(1)

    md_file = sys.argv[1]

    if not Path(md_file).exists():
        print(f"Error: File not found: {md_file}")
        sys.exit(1)

    # Determine output file
    if len(sys.argv) >= 3:
        output_file = sys.argv[2]
    else:
        output_file = Path(md_file).stem + '_with_images.pptx'
        output_file = Path('output/slides') / output_file

    # Ensure output directory exists
    Path(output_file).parent.mkdir(parents=True, exist_ok=True)

    # Convert
    print(f"📄 Converting {md_file} to PPTX (with image support)...")
    slides = parse_marpedown_with_images(md_file)
    print(f"   Found {len(slides)} slides")

    create_pptx_with_images(slides, str(output_file), title=Path(md_file).stem)


if __name__ == '__main__':
    main()
