#!/usr/bin/env python3
"""
Markdown to PPTX Converter

Converts Marp Markdown files to PowerPoint PPTX format
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re


def parse_marpedown(md_file: str):
    """Parse Marp Markdown file and extract slides"""

    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()

    # Remove YAML front matter
    if content.startswith('---'):
        parts = content.split('---', 2)
        if len(parts) >= 3:
            content = parts[2]

    # Split by slide separators
    slides_raw = content.split('\n---\n')

    slides = []
    for slide_raw in slides_raw:
        if not slide_raw.strip():
            continue

        # Parse slide
        lines = slide_raw.strip().split('\n')
        title = ""
        bullets = []

        for line in lines:
            line = line.strip()
            if not line:
                continue

            # Skip comments
            if line.startswith('<!--'):
                continue

            # Extract title (## Header)
            if line.startswith('## '):
                title = line[3:].strip()
            # Extract bullet points
            elif line.startswith('- ') or line.startswith('* '):
                bullets.append(line[2:].strip())

        if title or bullets:
            slides.append({
                'title': title,
                'bullets': bullets
            })

    return slides


def create_pptx(slides, output_file: str, title: str = "Presentation"):
    """Create PPTX from slides"""

    prs = Presentation()

    # Set slide dimensions to 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Create title slide
    if slides and slides[0]['title']:
        slide_layout = prs.slide_layouts[0]  # Title slide
        slide = prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = slides[0]['title']

        if slides[0]['bullets']:
            subtitle = slide.placeholders[1]
            subtitle.text = '\n'.join(slides[0]['bullets'][:2])

    # Create content slides
    for slide_data in slides[1:]:
        if not slide_data['title'] and not slide_data['bullets']:
            continue

        # Use blank layout
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Add title
        if slide_data['title']:
            left = Inches(0.5)
            top = Inches(0.5)
            width = Inches(12.333)
            height = Inches(1)

            title_box = slide.shapes.add_textbox(left, top, width, height)
            title_frame = title_box.text_frame
            title_frame.text = slide_data['title']

            # Style the title
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(36)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)  # Dark blue
            title_para.alignment = PP_ALIGN.LEFT

        # Add bullets
        if slide_data['bullets']:
            left = Inches(0.75)
            top = Inches(2)
            width = Inches(11.833)
            height = Inches(5)

            text_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = text_box.text_frame
            text_frame.word_wrap = True

            for i, bullet in enumerate(slide_data['bullets'][:6]):  # Max 6 bullets
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()

                p.text = bullet
                p.font.size = Pt(24)
                p.level = 0
                p.space_after = Pt(12)

                # Handle checkmarks and crosses
                if bullet.startswith('✓'):
                    p.font.color.rgb = RGBColor(0x27, 0xAE, 0x60)  # Green
                elif bullet.startswith('✗'):
                    p.font.color.rgb = RGBColor(0xE7, 0x4C, 0x3C)  # Red
                else:
                    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)  # Dark gray

    # Save presentation
    prs.save(output_file)
    print(f"✅ Created PPTX: {output_file}")
    print(f"   Total slides: {len(prs.slides)}")


def main():
    if len(sys.argv) < 2:
        print("Usage: python md_to_pptx.py <input.md> [output.pptx]")
        sys.exit(1)

    md_file = sys.argv[1]

    if not Path(md_file).exists():
        print(f"Error: File not found: {md_file}")
        sys.exit(1)

    # Determine output file
    if len(sys.argv) >= 3:
        output_file = sys.argv[2]
    else:
        output_file = Path(md_file).stem + '.pptx'
        output_file = Path('output/slides') / output_file

    # Ensure output directory exists
    Path(output_file).parent.mkdir(parents=True, exist_ok=True)

    # Convert
    print(f"📄 Converting {md_file} to PPTX...")
    slides = parse_marpedown(md_file)
    print(f"   Found {len(slides)} slides")

    create_pptx(slides, str(output_file), title=Path(md_file).stem)


if __name__ == '__main__':
    main()
